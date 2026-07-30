import os
import sys
try:
    import imagehash
    from PIL import Image
    from pptx import Presentation
except ImportError:
    os.system('pip install imagehash Pillow python-pptx')
    import imagehash
    from PIL import Image
    from pptx import Presentation
import io

pptx_path = r'C:\Users\Christian Gonzalez\Pictures\extras\Fichas tecnicas Aglael (1).pptx'
img_folder = r'C:\Users\Christian Gonzalez\Pictures\extras\CATALOGO DE OBRA AGLA\CATALOGO DE OBRA AGLA'

# Compute hashes for all original images
orig_hashes = {}
for file in os.listdir(img_folder):
    if file.endswith('.png') or file.endswith('.jpg'):
        path = os.path.join(img_folder, file)
        try:
            img = Image.open(path)
            orig_hashes[file] = imagehash.average_hash(img)
        except Exception as e:
            print(f"Error reading {file}: {e}")

prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    # Find first image in slide
    img_data = None
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            img_data = shape.image.blob
            break
    
    if img_data:
        try:
            slide_img = Image.open(io.BytesIO(img_data))
            slide_hash = imagehash.average_hash(slide_img)
            
            # Find closest
            best_match = None
            best_diff = float('inf')
            for name, h in orig_hashes.items():
                diff = slide_hash - h
                if diff < best_diff:
                    best_diff = diff
                    best_match = name
            
            print(f"Slide {slide_num} -> Best Match: {best_match} (diff: {best_diff})")
        except Exception as e:
            print(f"Slide {slide_num} -> Error processing image: {e}")
    else:
        print(f"Slide {slide_num} -> No image found in slide")
